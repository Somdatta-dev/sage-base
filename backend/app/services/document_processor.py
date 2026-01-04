"""
Document processing service using lightweight libraries.
Handles PDF, DOCX, PPTX, XLSX, and other document formats without GPU.
"""

import os
import tempfile
from pathlib import Path
from typing import Optional, Dict, Any
from io import BytesIO

# PDF processing
from pypdf import PdfReader

# Office document processing
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook

# HTML/Markdown processing
from bs4 import BeautifulSoup
import markdown


class DocumentProcessorService:
    """Service for processing documents using lightweight libraries."""
    
    def __init__(self):
        """Initialize the document processor."""
        pass
    
    async def process_document(
        self, 
        file_path: str,
    ) -> Dict[str, Any]:
        """
        Process a document and extract its content.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Dict containing:
            - markdown: Document content as markdown
            - text: Plain text content
            - metadata: Document metadata
        """
        ext = Path(file_path).suffix.lower()
        
        try:
            if ext == '.pdf':
                return await self._process_pdf(file_path)
            elif ext == '.docx':
                return await self._process_docx(file_path)
            elif ext == '.pptx':
                return await self._process_pptx(file_path)
            elif ext == '.xlsx':
                return await self._process_xlsx(file_path)
            elif ext == '.html':
                return await self._process_html(file_path)
            elif ext == '.md':
                return await self._process_markdown(file_path)
            elif ext == '.txt':
                return await self._process_text(file_path)
            else:
                return {
                    'success': False,
                    'error': f'Unsupported file type: {ext}',
                    'markdown': '',
                    'text': '',
                    'metadata': {},
                }
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'markdown': '',
                'text': '',
                'metadata': {},
            }
    
    async def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        """Process PDF file."""
        reader = PdfReader(file_path)
        
        text_parts = []
        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                text_parts.append(f"## Page {page_num}\n\n{text}")
        
        full_text = "\n\n".join(text_parts)
        
        # Get metadata
        metadata = {
            'filename': Path(file_path).name,
            'num_pages': len(reader.pages),
            'type': 'pdf',
        }
        
        if reader.metadata:
            if reader.metadata.title:
                metadata['title'] = reader.metadata.title
            if reader.metadata.author:
                metadata['author'] = reader.metadata.author
        
        return {
            'success': True,
            'markdown': full_text,
            'text': full_text,
            'metadata': metadata,
        }
    
    async def _process_docx(self, file_path: str) -> Dict[str, Any]:
        """Process DOCX file."""
        doc = DocxDocument(file_path)
        
        markdown_parts = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            # Detect headings by style
            if para.style and para.style.name:
                style_name = para.style.name.lower()
                if 'heading 1' in style_name:
                    markdown_parts.append(f"# {text}")
                elif 'heading 2' in style_name:
                    markdown_parts.append(f"## {text}")
                elif 'heading 3' in style_name:
                    markdown_parts.append(f"### {text}")
                elif 'title' in style_name:
                    markdown_parts.append(f"# {text}")
                else:
                    markdown_parts.append(text)
            else:
                markdown_parts.append(text)
        
        # Process tables
        for table in doc.tables:
            table_md = self._table_to_markdown(table)
            if table_md:
                markdown_parts.append(table_md)
        
        full_text = "\n\n".join(markdown_parts)
        
        return {
            'success': True,
            'markdown': full_text,
            'text': full_text,
            'metadata': {
                'filename': Path(file_path).name,
                'type': 'docx',
            },
        }
    
    def _table_to_markdown(self, table) -> str:
        """Convert a table to markdown format."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        
        if not rows:
            return ""
        
        # Create markdown table
        md_lines = []
        
        # Header
        md_lines.append("| " + " | ".join(rows[0]) + " |")
        md_lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
        
        # Body
        for row in rows[1:]:
            md_lines.append("| " + " | ".join(row) + " |")
        
        return "\n".join(md_lines)
    
    async def _process_pptx(self, file_path: str) -> Dict[str, Any]:
        """Process PPTX file."""
        prs = Presentation(file_path)
        
        markdown_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"## Slide {slide_num}"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if len(slide_text) > 1:  # More than just the slide header
                markdown_parts.append("\n\n".join(slide_text))
        
        full_text = "\n\n---\n\n".join(markdown_parts)
        
        return {
            'success': True,
            'markdown': full_text,
            'text': full_text,
            'metadata': {
                'filename': Path(file_path).name,
                'num_slides': len(prs.slides),
                'type': 'pptx',
            },
        }
    
    async def _process_xlsx(self, file_path: str) -> Dict[str, Any]:
        """Process XLSX file."""
        wb = load_workbook(file_path, read_only=True, data_only=True)
        
        markdown_parts = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            markdown_parts.append(f"## {sheet_name}")
            
            rows = []
            for row in sheet.iter_rows(values_only=True):
                # Convert None to empty string
                cells = [str(cell) if cell is not None else "" for cell in row]
                if any(cells):  # Skip completely empty rows
                    rows.append(cells)
            
            if rows:
                # Create markdown table
                md_lines = []
                max_cols = max(len(row) for row in rows)
                
                # Pad rows to same length
                for i, row in enumerate(rows):
                    rows[i] = row + [""] * (max_cols - len(row))
                
                # Header (first row)
                md_lines.append("| " + " | ".join(rows[0]) + " |")
                md_lines.append("| " + " | ".join(["---"] * max_cols) + " |")
                
                # Body
                for row in rows[1:]:
                    md_lines.append("| " + " | ".join(row) + " |")
                
                markdown_parts.append("\n".join(md_lines))
        
        wb.close()
        full_text = "\n\n".join(markdown_parts)
        
        return {
            'success': True,
            'markdown': full_text,
            'text': full_text,
            'metadata': {
                'filename': Path(file_path).name,
                'num_sheets': len(wb.sheetnames),
                'type': 'xlsx',
            },
        }
    
    async def _process_html(self, file_path: str) -> Dict[str, Any]:
        """Process HTML file."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Get text
        text = soup.get_text(separator='\n')
        
        # Clean up whitespace
        lines = [line.strip() for line in text.splitlines()]
        text = '\n'.join(line for line in lines if line)
        
        # Get title if available
        title = soup.title.string if soup.title else None
        
        return {
            'success': True,
            'markdown': text,
            'text': text,
            'metadata': {
                'filename': Path(file_path).name,
                'title': title,
                'type': 'html',
            },
        }
    
    async def _process_markdown(self, file_path: str) -> Dict[str, Any]:
        """Process Markdown file."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            md_content = f.read()
        
        # Convert to HTML and back to plain text for text field
        html = markdown.markdown(md_content)
        soup = BeautifulSoup(html, 'html.parser')
        plain_text = soup.get_text()
        
        return {
            'success': True,
            'markdown': md_content,
            'text': plain_text,
            'metadata': {
                'filename': Path(file_path).name,
                'type': 'markdown',
            },
        }
    
    async def _process_text(self, file_path: str) -> Dict[str, Any]:
        """Process plain text file."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text_content = f.read()
        
        return {
            'success': True,
            'markdown': text_content,
            'text': text_content,
            'metadata': {
                'filename': Path(file_path).name,
                'type': 'text',
            },
        }
    
    async def process_uploaded_file(
        self,
        file_content: bytes,
        filename: str,
    ) -> Dict[str, Any]:
        """
        Process an uploaded file from bytes.
        
        Args:
            file_content: File content as bytes
            filename: Original filename
            
        Returns:
            Processing result dict
        """
        # Create a temporary file
        suffix = Path(filename).suffix
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp_file:
            tmp_file.write(file_content)
            tmp_path = tmp_file.name
        
        try:
            result = await self.process_document(tmp_path)
            result['metadata']['original_filename'] = filename
            return result
        finally:
            # Clean up temp file
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
    
    def _parse_inline_formatting(self, text: str) -> list:
        """
        Parse inline markdown formatting into TipTap content array.

        Supports: [text](url) links, plain URLs, **bold**, *italic*, `code`, ~~strikethrough~~, ==highlight==

        Args:
            text: Text with markdown inline formatting

        Returns:
            List of TipTap text nodes with marks
        """
        import re

        # Pattern to match all inline formatting
        # Order matters: links, URLs, code, highlight, strikethrough, bold, italic
        result = []
        current_pos = 0

        # Comprehensive pattern for all inline formatting
        pattern = r'(\[([^\]]+)\]\(([^)]+)\))|(https?://[^\s<>)]+|www\.[^\s<>)]+)|(`[^`]+`)|(==[^=]+==)|(~~[^~]+~~)|(\*\*[^*]+\*\*)|(\*[^*]+\*)'

        for match in re.finditer(pattern, text):
            # Add text before the match (plain text)
            if match.start() > current_pos:
                plain = text[current_pos:match.start()]
                if plain:
                    result.append({'type': 'text', 'text': plain})

            matched_text = match.group()

            # Link: [text](url)
            if matched_text.startswith('[') and '](' in matched_text:
                link_text = match.group(2)
                link_url = match.group(3)
                result.append({
                    'type': 'text',
                    'text': link_text,
                    'marks': [{'type': 'link', 'attrs': {'href': link_url, 'target': '_blank'}}]
                })
            # Auto-link plain URLs
            elif matched_text.startswith('http') or matched_text.startswith('www'):
                url = matched_text if matched_text.startswith('http') else f'https://{matched_text}'
                result.append({
                    'type': 'text',
                    'text': matched_text,
                    'marks': [{'type': 'link', 'attrs': {'href': url, 'target': '_blank'}}]
                })
            # Inline code
            elif matched_text.startswith('`') and matched_text.endswith('`'):
                result.append({
                    'type': 'text',
                    'text': matched_text[1:-1],
                    'marks': [{'type': 'code'}]
                })
            # Highlight: ==text==
            elif matched_text.startswith('==') and matched_text.endswith('=='):
                result.append({
                    'type': 'text',
                    'text': matched_text[2:-2],
                    'marks': [{'type': 'highlight'}]
                })
            # Strikethrough: ~~text~~
            elif matched_text.startswith('~~') and matched_text.endswith('~~'):
                result.append({
                    'type': 'text',
                    'text': matched_text[2:-2],
                    'marks': [{'type': 'strike'}]
                })
            # Bold
            elif matched_text.startswith('**') and matched_text.endswith('**'):
                result.append({
                    'type': 'text',
                    'text': matched_text[2:-2],
                    'marks': [{'type': 'bold'}]
                })
            # Italic
            elif matched_text.startswith('*') and matched_text.endswith('*'):
                result.append({
                    'type': 'text',
                    'text': matched_text[1:-1],
                    'marks': [{'type': 'italic'}]
                })

            current_pos = match.end()

        # Add remaining text
        if current_pos < len(text):
            remaining = text[current_pos:]
            if remaining:
                result.append({'type': 'text', 'text': remaining})

        # If no formatting found, return simple text node
        if not result:
            return [{'type': 'text', 'text': text}] if text else []

        return result

    def _parse_list_items(self, lines: list, start_index: int, list_type: str = 'bullet') -> tuple[list, int]:
        """
        Parse list items with support for nested lists.

        Args:
            lines: All lines of the document
            start_index: Index where list starts
            list_type: 'bullet' or 'ordered'

        Returns:
            Tuple of (list of list items, next line index after list)
        """
        list_items = []
        i = start_index

        def get_indent_level(line: str) -> int:
            """Get the indentation level (number of spaces/tabs at start)"""
            return len(line) - len(line.lstrip())

        def is_list_item(line: str, ltype: str) -> bool:
            """Check if line is a list item of the given type"""
            stripped = line.lstrip()
            if ltype == 'bullet':
                return stripped.startswith('- ') or stripped.startswith('* ')
            else:  # ordered
                return stripped and stripped[0].isdigit() and '. ' in stripped[:4]

        base_indent = get_indent_level(lines[i]) if i < len(lines) else 0

        while i < len(lines):
            line = lines[i]
            if not line.strip():
                i += 1
                continue

            indent = get_indent_level(line)

            # If line is not indented properly for this list level, we're done
            if indent < base_indent:
                break

            # Check if it's a list item at current level
            if indent == base_indent and is_list_item(line, list_type):
                # Extract item text
                stripped = line.lstrip()
                if list_type == 'bullet':
                    item_text = stripped[2:].strip()
                else:
                    item_text = stripped.split('. ', 1)[1].strip() if '. ' in stripped else stripped

                item_content = [{
                    'type': 'paragraph',
                    'content': self._parse_inline_formatting(item_text)
                }]

                # Check for nested list on next line
                if i + 1 < len(lines):
                    next_line = lines[i + 1]
                    next_indent = get_indent_level(next_line)

                    if next_indent > indent and next_line.strip():
                        # Nested list detected
                        nested_type = 'bullet' if (next_line.lstrip().startswith('- ') or next_line.lstrip().startswith('* ')) else 'ordered'
                        nested_items, next_i = self._parse_list_items(lines, i + 1, nested_type)

                        if nested_items:
                            item_content.append({
                                'type': 'bulletList' if nested_type == 'bullet' else 'orderedList',
                                'content': nested_items
                            })
                            i = next_i
                            list_items.append({
                                'type': 'listItem',
                                'content': item_content
                            })
                            continue

                list_items.append({
                    'type': 'listItem',
                    'content': item_content
                })
                i += 1
            elif indent > base_indent:
                # Skip this line, it will be handled by nested list parsing
                i += 1
            else:
                # Not a list item at this level, we're done
                break

        return list_items, i

    def _parse_table(self, lines: list, start_index: int) -> tuple[Dict[str, Any], int]:
        """
        Parse a markdown table into TipTap table JSON structure.

        Args:
            lines: All lines of the document
            start_index: Index where table starts

        Returns:
            Tuple of (table JSON structure, next line index after table)
        """
        table_lines = []
        i = start_index

        # Collect all table lines
        while i < len(lines) and lines[i].startswith('|'):
            table_lines.append(lines[i])
            i += 1

        if not table_lines:
            return None, start_index

        # Parse table structure
        rows = []
        is_first_row = True
        has_header = False

        for line_idx, line in enumerate(table_lines):
            # Check if this is the separator row (|---|---|)
            if '---' in line or ':-:' in line or ':--' in line or '--:' in line:
                has_header = True
                continue

            # Parse cells from the row
            cells = [c.strip() for c in line.split('|')]
            # Remove empty first/last cells from | cell | cell | format
            cells = [c for c in cells if c]

            if not cells:
                continue

            # Create table cells with inline formatting
            cell_nodes = []
            for cell_text in cells:
                cell_content = self._parse_inline_formatting(cell_text)
                cell_nodes.append({
                    'type': 'tableCell' if not (is_first_row and has_header) or line_idx > 0 else 'tableHeader',
                    'content': [{
                        'type': 'paragraph',
                        'content': cell_content if cell_content else [{'type': 'text', 'text': ''}]
                    }]
                })

            rows.append({
                'type': 'tableRow',
                'content': cell_nodes
            })

            if is_first_row and has_header:
                is_first_row = False

        if not rows:
            return None, start_index

        return {
            'type': 'table',
            'content': rows
        }, i

    def convert_to_tiptap_json(self, markdown_text: str) -> Dict[str, Any]:
        """
        Convert markdown content to Tiptap-compatible JSON.

        This creates a structured document that preserves:
        - Headings
        - Paragraphs
        - Code blocks
        - Lists
        - Tables (as HTML)
        - Inline formatting (bold, italic, code)
        """
        lines = markdown_text.split('\n')
        content = []
        current_code_block = None
        code_language = 'text'
        
        i = 0
        while i < len(lines):
            line = lines[i]
            
            # Handle code blocks
            if line.startswith('```'):
                if current_code_block is None:
                    # Start of code block
                    code_language = line[3:].strip() or 'text'
                    current_code_block = []
                else:
                    # End of code block
                    content.append({
                        'type': 'codeBlock',
                        'attrs': {'language': code_language},
                        'content': [{'type': 'text', 'text': '\n'.join(current_code_block)}]
                    })
                    current_code_block = None
                i += 1
                continue
            
            if current_code_block is not None:
                current_code_block.append(line)
                i += 1
                continue
            
            # Handle headings
            if line.startswith('# '):
                content.append({
                    'type': 'heading',
                    'attrs': {'level': 1},
                    'content': self._parse_inline_formatting(line[2:].strip())
                })
            elif line.startswith('## '):
                content.append({
                    'type': 'heading',
                    'attrs': {'level': 2},
                    'content': self._parse_inline_formatting(line[3:].strip())
                })
            elif line.startswith('### '):
                content.append({
                    'type': 'heading',
                    'attrs': {'level': 3},
                    'content': self._parse_inline_formatting(line[4:].strip())
                })
            # Handle bullet lists (with nested list support)
            elif line.startswith('- ') or line.startswith('* '):
                list_items, next_i = self._parse_list_items(lines, i, 'bullet')
                if list_items:
                    content.append({
                        'type': 'bulletList',
                        'content': list_items
                    })
                    i = next_i
                    continue
                else:
                    i += 1
                    continue
            # Handle numbered lists (with nested list support)
            elif line and line[0].isdigit() and '. ' in line[:4]:
                list_items, next_i = self._parse_list_items(lines, i, 'ordered')
                if list_items:
                    content.append({
                        'type': 'orderedList',
                        'content': list_items
                    })
                    i = next_i
                    continue
                else:
                    i += 1
                    continue
            # Handle blockquotes
            elif line.startswith('> '):
                quote_lines = []
                while i < len(lines) and lines[i].startswith('> '):
                    quote_lines.append(lines[i][2:])
                    i += 1
                quote_text = '\n'.join(quote_lines)
                content.append({
                    'type': 'blockquote',
                    'content': [{
                        'type': 'paragraph',
                        'content': self._parse_inline_formatting(quote_text)
                    }]
                })
                continue
            # Handle horizontal rules
            elif line.strip() in ['---', '***', '___']:
                content.append({'type': 'horizontalRule'})
            # Handle table rows (markdown tables)
            elif line.startswith('|'):
                # Parse table using helper function
                table_node, next_i = self._parse_table(lines, i)
                if table_node:
                    content.append(table_node)
                    i = next_i
                    continue
                else:
                    # Fallback: treat as plain text if parsing fails
                    i += 1
                    continue
            # Handle regular paragraphs
            elif line.strip():
                content.append({
                    'type': 'paragraph',
                    'content': self._parse_inline_formatting(line)
                })
            
            i += 1
        
        return {
            'type': 'doc',
            'content': content if content else [{'type': 'paragraph', 'content': []}]
        }


# Singleton instance
_document_processor: Optional[DocumentProcessorService] = None


def get_document_processor() -> DocumentProcessorService:
    """Get or create the document processor singleton."""
    global _document_processor
    if _document_processor is None:
        _document_processor = DocumentProcessorService()
    return _document_processor
